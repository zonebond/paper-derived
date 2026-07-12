"""统一文件格式读取层.

检测文件格式 → 提取纯文本内容 → 返回 (text, format_type, metadata).

引擎永远只看到纯文本, 格式转换完全在 CLI 层完成。
"""

from __future__ import annotations

import csv
import io
import os
import re
import shutil
import subprocess
import sys
import tempfile
from collections import defaultdict
from html.parser import HTMLParser
from pathlib import Path


# ── 扩展名 → 格式类型映射 ──────────────────────────────────────

SUPPORTED_EXTENSIONS: dict[str, str] = {
    ".docx": "docx",
    ".doc": "doc",
    ".xlsx": "xlsx",
    ".xls": "xls",
    ".pdf": "pdf",
    ".pptx": "pptx",
    ".csv": "csv",
    ".tsv": "tsv",
    ".md": "markdown",
    ".markdown": "markdown",
    ".txt": "plain_text",
    ".text": "plain_text",
    ".json": "json",
}

# 纯文本格式 — 直接用 Path.read_text() 读取
_TEXT_FORMATS = {".md", ".markdown", ".txt", ".text", ".json"}


def detect_format(file_path: str | Path) -> str:
    """由文件扩展名推断格式类型.

    Args:
        file_path: 文件路径.

    Returns:
        格式类型字符串 (docx, xlsx, pdf, markdown, plain_text, json, ...).

    Raises:
        ValueError: 不支持的扩展名.
    """
    suffix = Path(file_path).suffix.lower()
    fmt = SUPPORTED_EXTENSIONS.get(suffix)
    if fmt is None:
        raise ValueError(
            f"不支持的文件格式: {suffix}。"
            f"支持的格式: {', '.join(sorted(SUPPORTED_EXTENSIONS.keys()))}"
        )
    return fmt


def read_file(file_path: str | Path) -> tuple[str, str, dict]:
    """读取任意支持格式的文件, 返回纯文本内容.

    Args:
        file_path: 文件路径.

    Returns:
        (content: str, format_type: str, metadata: dict)
        - content: 提取的纯文本内容
        - format_type: 格式类型字符串
        - metadata: 附加元数据 (页数、表数等)

    Raises:
        ValueError: 不支持的格式.
        FileNotFoundError: 文件不存在.
        RuntimeError: 读取失败 (含依赖缺失提示).
    """
    path = Path(file_path).expanduser().resolve()

    if not path.exists():
        raise FileNotFoundError(f"文件不存在: {path}")

    fmt = detect_format(path)
    metadata: dict = {"source": str(path), "format": fmt, "size_bytes": path.stat().st_size}

    # 纯文本 - 直接读取
    if path.suffix.lower() in _TEXT_FORMATS:
        content = path.read_text(encoding="utf-8")
        return content, fmt, metadata

    # 二进制/结构化格式
    reader = _READERS.get(fmt)
    if reader is None:
        raise RuntimeError(f"格式 '{fmt}' 暂不支持读取")

    content, extra_meta = reader(path)
    metadata.update(extra_meta)
    return content, fmt, metadata


# ── 各格式读取器 ──────────────────────────────────────────────


def _read_docx(path: Path) -> tuple[str, dict]:
    """读取 .docx 文件, 提取段落文本和表格."""
    try:
        from docx import Document
    except ImportError:
        raise RuntimeError("读取 .docx 需要 python-docx: pip install python-docx")

    doc = Document(path)
    parts: list[str] = []
    table_count = 0

    # 迭代文档 body 中的所有块元素 (段落 + 表格), 保持文档顺序
    body = doc.element.body
    for child in body:
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag

        if tag == "p":
            # 段落: 在 doc.paragraphs 中找到对应的 Paragraph 对象
            for para in doc.paragraphs:
                if para._element is child:
                    text = para.text.strip()
                    if text:
                        style_name = para.style.name if para.style else ""
                        if style_name and style_name.startswith("Heading"):
                            try:
                                level = int(style_name.split()[-1])
                            except ValueError:
                                level = 1
                            parts.append(f"{'#' * max(1, min(6, level))} {text}")
                        else:
                            parts.append(text)
                    break

        elif tag == "tbl":
            table_count += 1
            for table in doc.tables:
                if table._element is child:
                    parts.append(_table_to_text(table))
                    break
            else:
                parts.append("[表格]")

    return "\n\n".join(parts), {"table_count": table_count}


def _read_doc(path: Path) -> tuple[str, dict]:
    """读取旧版 .doc 文件 (Composite Document File V2).

    优先通过 textutil HTML 转换保留 Word 标题层级（Heading 1/2/3 → #/##/###）。
    若 HTML 解析失败则回退到纯文本转换。

    macOS: 使用内置 textutil.
    Linux: 尝试 antiword 或 catdoc。
    """
    # macOS — textutil with HTML-based heading preservation
    if sys.platform == "darwin":
        if not shutil.which("textutil"):
            raise RuntimeError("macOS 上需要 textutil (系统内置)")
        try:
            # Try HTML conversion first (preserves heading styles)
            annotated = _doc_via_html(path)
            if annotated is not None:
                return annotated, {"converter": "textutil+html"}
            # Fallback: plain text
            content = _doc_via_textutil_txt(path)
            return content, {"converter": "textutil"}
        except subprocess.TimeoutExpired:
            raise RuntimeError("textutil 转换超时")
        except FileNotFoundError:
            raise RuntimeError("textutil 不可用")

    # Linux — antiword
    if shutil.which("antiword"):
        try:
            result = subprocess.run(
                ["antiword", str(path)], capture_output=True, text=True, timeout=30
            )
            if result.returncode == 0:
                return result.stdout, {"converter": "antiword"}
        except (subprocess.TimeoutExpired, FileNotFoundError):
            pass

    # Linux — catdoc
    if shutil.which("catdoc"):
        try:
            result = subprocess.run(
                ["catdoc", str(path)], capture_output=True, text=True, timeout=30
            )
            if result.returncode == 0:
                return result.stdout, {"converter": "catdoc"}
        except (subprocess.TimeoutExpired, FileNotFoundError):
            pass

    raise RuntimeError(
        "无法读取 .doc 文件。"
        "macOS: 系统自带 textutil 应可用。"
        "Linux: 请安装 antiword (apt install antiword) 或 catdoc (apt install catdoc)。"
    )


def _doc_via_textutil_txt(path: Path) -> str:
    """Convert .doc to plain text via textutil (no heading preservation)."""
    with tempfile.NamedTemporaryFile(suffix=".txt", delete=False) as tmp:
        tmp_path = tmp.name
    try:
        result = subprocess.run(
            ["textutil", "-convert", "txt", str(path), "-output", tmp_path],
            capture_output=True, text=True, timeout=30,
        )
        if result.returncode != 0:
            raise RuntimeError(f"textutil 转换失败: {result.stderr}")
        return Path(tmp_path).read_text(encoding="utf-8")
    finally:
        os.unlink(tmp_path)


# ── .doc HTML → 标注文本（保留 Word 标题层级）─────────────────────


# Regex for TOC entries: "HYPERLINK \l "_TocNNN" N[.N] Title\t PAGEREF ..."
_TOC_LINE_RE = re.compile(
    r'HYPERLINK\s+\\l\s+"_[^"]*"\s*'
    r'(\d+(?:\.\d+)*)\s+'
    r'(.+?)'
    r'\t\s*PAGEREF'
)

_CAPTION_RE = re.compile(r'^(表|图)\s*\d+[-–—]\d+')
_LABEL_RE = re.compile(r'^[（(][^）)]+[）)]$')


class _DocHTMLParser(HTMLParser):
    """Parse textutil HTML output, collecting paragraphs and CSS."""

    def __init__(self):
        super().__init__()
        self.style_block = ""
        self.paragraphs = []  # [{class, text, bold}, ...]

        self._in_style = False
        self._style_buf = []
        self._current_class = ""
        self._buffer = []
        self._had_bold = False
        self._in_bold = False

    def handle_starttag(self, tag, attrs):
        d = dict(attrs)
        if tag == 'style':
            self._in_style = True
            self._style_buf = []
        elif tag == 'p':
            self._current_class = d.get('class', '')
            self._buffer = []
            self._had_bold = False
        elif tag == 'b':
            self._in_bold = True

    def handle_endtag(self, tag):
        if tag == 'style':
            self._in_style = False
            self.style_block = ''.join(self._style_buf)
        elif tag == 'p':
            text = ''.join(self._buffer).strip()
            if text:
                self.paragraphs.append({
                    'class': self._current_class,
                    'text': text,
                    'bold': self._had_bold,
                })
            self._current_class = ""
            self._buffer = []
        elif tag == 'b':
            self._in_bold = False

    def handle_data(self, data):
        if self._in_style:
            self._style_buf.append(data)
        elif self._current_class:
            if self._in_bold:
                self._had_bold = True
            self._buffer.append(data)


def _parse_css(style_block: str) -> dict:
    """Parse <style> → {class_name: {font_size, margin_left, line_height, ...}}."""
    props = {}
    for m in re.finditer(r'p\.(p\d+)\s*\{([^}]+)\}', style_block, re.DOTALL):
        cls = m.group(1)
        body = m.group(2)

        fm = re.search(r'font:\s*([\d.]+)px', body)
        font_size = float(fm.group(1)) if fm else None

        mlm = re.search(r'margin-left:\s*([\d.-]+)px', body)
        if mlm:
            margin_left = float(mlm.group(1))
        else:
            sh = re.search(r'margin:\s*[\d.]+\w*\s+[\d.]+\w*\s+[\d.]+\w*\s+([\d.]+)\w*', body)
            margin_left = float(sh.group(1)) if sh else 0.0

        mtm = re.search(r'margin-top:\s*([\d.]+)px', body)
        if mtm:
            margin_top = float(mtm.group(1))
        else:
            sh = re.search(r'margin:\s*([\d.]+)\w*', body)
            margin_top = float(sh.group(1)) if sh else 0.0

        lhm = re.search(r'line-height:\s*([\d.]+)px', body)
        line_height = float(lhm.group(1)) if lhm else None

        tam = re.search(r'text-align:\s*(\w+)', body)
        text_align = tam.group(1) if tam else 'justify'

        props[cls] = {
            'font_size': font_size,
            'margin_left': margin_left,
            'margin_top': margin_top,
            'line_height': line_height,
            'text_align': text_align,
        }
    return props


def _normalize(text: str) -> str:
    """Normalize for fuzzy matching: strip whitespace and HTML entities."""
    t = text.strip()
    t = re.sub(r'\s+', '', t)
    t = t.replace('&amp;', '&')
    return t


def _doc_via_html(path: Path) -> str | None:
    """Convert .doc → heading-annotated text via HTML parsing.

    Strategy (fully data-driven):
      1. textutil -convert html → parse HTML
      2. Extract TOC entries (HYPERLINK+PAGEREF) → ground-truth heading levels
      3. Extract CSS properties per class
      4. Match TOC titles to body bold paragraphs → CSS class → level map
      5. For bold paragraphs not in TOC, infer level from margin-left:
         if ml > max_known_ml + 10 → deeper level
         else → CSS similarity to nearest known class
      6. Annotate plain text with #/##/### markers

    Returns annotated text, or None if HTML parsing fails.
    """
    # Step 1: Convert .doc → HTML
    html = None
    with tempfile.NamedTemporaryFile(suffix='.html', delete=False) as tmp:
        html_path = tmp.name
    try:
        result = subprocess.run(
            ['textutil', '-convert', 'html', str(path), '-output', html_path],
            capture_output=True, text=True, timeout=30,
        )
        if result.returncode != 0:
            return None
        html = Path(html_path).read_text(encoding='utf-8')
    except Exception:
        return None
    finally:
        os.unlink(html_path)

    if not html:
        return None

    # Step 2: Parse HTML
    parser = _DocHTMLParser()
    try:
        parser.feed(html)
    except Exception:
        return None

    # Step 3: Separate TOC from body
    toc_entries = []  # [(num_str, title, level), ...]
    body_paras = []
    non_toc_streak = 0
    for p in parser.paragraphs:
        if 'HYPERLINK' in p['text'] and 'PAGEREF' in p['text']:
            m = _TOC_LINE_RE.search(p['text'])
            if m:
                num_str = m.group(1)
                title = _normalize(m.group(2))
                level = num_str.count('.') + 1
                toc_entries.append((num_str, title, level))
                non_toc_streak = 0
                continue
        non_toc_streak += 1
        if non_toc_streak >= 2:
            body_paras.append(p)

    if not toc_entries or not body_paras:
        return None

    # Step 4: Extract CSS
    css = _parse_css(parser.style_block)

    # Step 5: Match TOC → body bold paragraphs → class → level
    bold_by_norm = defaultdict(list)
    for p in body_paras:
        if p['bold']:
            bold_by_norm[_normalize(p['text'])].append(p)

    class_votes = defaultdict(lambda: defaultdict(int))
    heading_map = {}  # text → (0based_level, class)

    for _num, toc_title, toc_level in toc_entries:
        candidates = bold_by_norm.get(toc_title, [])
        if not candidates:
            # Fuzzy match
            for norm, paras in bold_by_norm.items():
                if len(toc_title) >= 4 and (toc_title in norm or norm in toc_title):
                    candidates = paras
                    break
        if candidates:
            cls = candidates[0]['class']
            class_votes[cls][toc_level - 1] += 1
            heading_map[candidates[0]['text']] = (toc_level - 1, cls)

    # Determine class → level
    class_level = {}
    for cls, votes in class_votes.items():
        class_level[cls] = max(votes, key=votes.get)

    if not class_level:
        return None

    # Step 6: Infer levels for unknown bold classes
    max_known_ml = 0.0
    for cls in class_level:
        if cls in css:
            ml = css[cls].get('margin_left', 0.0)
            if ml > max_known_ml:
                max_known_ml = ml

    max_known_level = max(class_level.values())

    for p in body_paras:
        if not p['bold']:
            continue
        cls = p['class']
        text = p['text']
        if cls in class_level:
            if text not in heading_map and not _CAPTION_RE.match(text) and not _LABEL_RE.match(text):
                heading_map[text] = (class_level[cls], cls)
        elif cls in css:
            feats = css[cls]
            if feats.get('text_align') == 'center':
                continue
            fs = feats.get('font_size')
            if fs is not None and fs <= 10.5:
                continue
            ml = feats.get('margin_left', 0.0)
            if ml > max_known_ml + 10 or ml > max_known_ml * 1.5:
                class_level[cls] = max_known_level + 1
            else:
                # CSS similarity match
                best_cls, best_dist = None, float('inf')
                for kcls, _klv in class_level.items():
                    if kcls not in css:
                        continue
                    kf = css[kcls]
                    d = abs(feats.get('font_size', 14) - kf.get('font_size', 14)) * 10 \
                        + abs(ml - kf.get('margin_left', 0)) \
                        + abs(feats.get('line_height', 24) - kf.get('line_height', 24)) * 2
                    if d < best_dist:
                        best_dist = d
                        best_cls = kcls
                if best_cls and best_dist < 60:
                    class_level[cls] = class_level[best_cls]
            if cls in class_level and text not in heading_map \
                    and not _CAPTION_RE.match(text) and not _LABEL_RE.match(text):
                heading_map[text] = (class_level[cls], cls)

    if not heading_map:
        return None

    # Step 7: Annotate plain text
    plain_text = _doc_via_textutil_txt(path)
    lines = plain_text.split('\n')
    annotated = []
    for line in lines:
        s = line.strip()
        if s in heading_map:
            lv, _ = heading_map[s]
            prefix = '#' * max(lv + 1, 1) + ' '
            annotated.append(f"{prefix}{s}")
        else:
            annotated.append(line)
    return '\n'.join(annotated)


def _read_xlsx(path: Path) -> tuple[str, dict]:
    """读取 .xlsx 文件, 保留表格结构."""
    try:
        import openpyxl
    except ImportError:
        raise RuntimeError("读取 .xlsx 需要 openpyxl: pip install openpyxl")

    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    parts: list[str] = []
    sheet_count = 0

    for sheet_name in wb.sheetnames:
        sheet_count += 1
        ws = wb[sheet_name]
        parts.append(f"## Sheet: {sheet_name}")

        # 读取最多 500 行，避免超大文件
        rows = list(ws.iter_rows(max_row=500, values_only=True))
        if not rows:
            parts.append("(空工作表)")
            continue

        for row in rows:
            cells = [str(c) if c is not None else "" for c in row]
            parts.append("| " + " | ".join(cells) + " |")

        if ws.max_row and ws.max_row > 500:
            parts.append(f"\n*(截断: 仅显示前 500 行, 共 {ws.max_row} 行)*")

    wb.close()
    return "\n\n".join(parts), {"sheet_count": sheet_count}


def _read_xls(path: Path) -> tuple[str, dict]:
    """读取旧版 .xls 文件."""
    try:
        import xlrd
    except ImportError:
        raise RuntimeError("读取 .xls 需要 xlrd: pip install xlrd")

    wb = xlrd.open_workbook(str(path))
    parts: list[str] = []
    sheet_count = 0

    for sheet_name in wb.sheet_names():
        sheet_count += 1
        ws = wb.sheet_by_name(sheet_name)
        parts.append(f"## Sheet: {sheet_name}")

        max_rows = min(ws.nrows, 500)
        for r in range(max_rows):
            cells = [str(ws.cell_value(r, c)) if ws.cell_value(r, c) != "" else ""
                     for c in range(ws.ncols)]
            parts.append("| " + " | ".join(cells) + " |")

        if ws.nrows > 500:
            parts.append(f"\n*(截断: 仅显示前 500 行, 共 {ws.nrows} 行)*")

    return "\n\n".join(parts), {"sheet_count": sheet_count}


def _read_pdf(path: Path) -> tuple[str, dict]:
    """读取 PDF 文件, 按页提取文本."""
    try:
        from pypdf import PdfReader
    except ImportError:
        raise RuntimeError("读取 PDF 需要 pypdf: pip install pypdf")

    reader = PdfReader(str(path))
    parts: list[str] = []
    page_count = len(reader.pages)

    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text()
        if text and text.strip():
            parts.append(f"--- Page {i} ---\n{text.strip()}")

    return "\n\n".join(parts), {"page_count": page_count}


def _read_pptx(path: Path) -> tuple[str, dict]:
    """读取 .pptx 文件, 提取每张幻灯片的文本."""
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("读取 .pptx 需要 python-pptx: pip install python-pptx")

    prs = Presentation(str(path))
    parts: list[str] = []
    slide_count = len(prs.slides)

    for i, slide in enumerate(prs.slides, start=1):
        parts.append(f"## Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        parts.append(para.text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text for cell in row.cells]
                    parts.append("| " + " | ".join(cells) + " |")

    return "\n\n".join(parts), {"slide_count": slide_count}


def _read_csv(path: Path) -> tuple[str, dict]:
    """读取 CSV 文件, 自动检测分隔符."""
    text = path.read_text(encoding="utf-8")

    # 自动检测分隔符: 统计每行中逗号 vs 制表符
    sample = text[:4096]
    comma_count = sample.count(",")
    tab_count = sample.count("\t")
    delimiter = "\t" if tab_count > comma_count else ","

    reader = csv.reader(io.StringIO(text), delimiter=delimiter)
    rows = list(reader)
    row_count = len(rows)

    parts = []
    for row in rows[:500]:
        parts.append("| " + " | ".join(row) + " |")

    if row_count > 500:
        parts.append(f"\n*(截断: 仅显示前 500 行, 共 {row_count} 行)*")

    return "\n".join(parts), {"row_count": row_count, "delimiter": delimiter}


# ── 读取器注册表 ──────────────────────────────────────────────

_READERS: dict[str, callable] = {
    "docx": _read_docx,
    "doc": _read_doc,
    "xlsx": _read_xlsx,
    "xls": _read_xls,
    "pdf": _read_pdf,
    "pptx": _read_pptx,
    "csv": _read_csv,
    "tsv": _read_csv,  # TSV 复用 CSV 读取器 (自动检测制表符)
}


# ── 分块 ───────────────────────────────────────────────────────


# 中文约 1.5 token/字，英文约 1 token/4 chars，取保守值
# 30000 chars ≈ 10K–15K tokens，留足空间给 system prompt
DEFAULT_CHUNK_SIZE = 30_000


def chunk_text(text: str, max_chars: int = DEFAULT_CHUNK_SIZE) -> list[str]:
    """将长文本按自然边界分块.

    优先在章节标题（# 标记）处切割，其次在空行处切割。
    每块不超过 max_chars 字符。

    Args:
        text: 原始文本.
        max_chars: 每块最大字符数.

    Returns:
        文本块列表（至少 1 块）.
    """
    if len(text) <= max_chars:
        return [text]

    # 按「连续两个换行」拆成段，保留分隔符用于还原
    paragraphs = text.split("\n\n")
    chunks: list[str] = []
    current: list[str] = []
    current_len = 0

    for para in paragraphs:
        para_len = len(para) + 2  # +2 for \n\n separator

        # 单段超长 → 强制按行拆
        if para_len > max_chars:
            if current:
                chunks.append("\n\n".join(current))
                current, current_len = [], 0
            # 按行拆分超长段
            for line in para.split("\n"):
                line_len = len(line) + 1
                if current_len + line_len > max_chars and current:
                    chunks.append("\n".join(current))
                    current, current_len = [], 0
                current.append(line)
                current_len += line_len
            continue

        # 当前块加这段就超了 → 先封块
        if current_len + para_len > max_chars and current:
            chunks.append("\n\n".join(current))
            current, current_len = [], 0

        current.append(para)
        current_len += para_len

    if current:
        chunks.append("\n\n".join(current))

    return chunks if chunks else [text]


# ── 辅助函数 ─────────────────────────────────────────────────


def _table_to_text(table) -> str:
    """将 python-docx 表格转为 Markdown 表格文本."""
    lines = []
    for row in table.rows:
        cells = [cell.text.replace("\n", " ") for cell in row.cells]
        lines.append("| " + " | ".join(cells) + " |")
    return "\n".join(lines)
